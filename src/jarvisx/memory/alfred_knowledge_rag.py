"""
Alfred Knowledge RAG & Semantic Vector Indexer.
Dedicated to Charan's personal engineering, coding, DSA, and academic subjects.
Indexes:
  1. Computer System Architecture & Operating Systems Lecture PPTs
  2. Data Structures & Algorithms (DSA) Code & Tutorials
  3. Numpy & Matrix Mathematics Computing Modules
  4. Friday 10-CGPA Academic Milestones & Governance Docs
"""

from __future__ import annotations

import json
import logging
import os
import sqlite3
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

logger = logging.getLogger("alfred_knowledge_rag")


@dataclass
class KnowledgeDocument:
    doc_id: str
    title: str
    category: str
    content: str
    source_path: str
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class SearchResult:
    doc_id: str
    title: str
    category: str
    snippet: str
    score: float
    source_path: str
    metadata: Dict[str, Any] = field(default_factory=dict)


class AlfredKnowledgeRAG:
    """Offline Semantic Vector RAG Knowledge Base for Alfred."""

    _instance: Optional[AlfredKnowledgeRAG] = None

    def __init__(self, db_path: str = "var/db/alfred_knowledge_rag.db"):
        self.db_path = db_path
        os.makedirs(os.path.dirname(self.db_path), exist_ok=True)
        self.vectorizer: Optional[TfidfVectorizer] = None
        self.tfidf_matrix: Optional[np.ndarray] = None
        self.documents: List[KnowledgeDocument] = []
        self._init_db()
        self._load_or_index_knowledge()

    @classmethod
    def get_instance(cls) -> AlfredKnowledgeRAG:
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance

    def _init_db(self):
        with sqlite3.connect(self.db_path) as conn:
            conn.execute("""
                CREATE TABLE IF NOT EXISTS alfred_knowledge_docs (
                    doc_id TEXT PRIMARY KEY,
                    title TEXT NOT NULL,
                    category TEXT NOT NULL,
                    content TEXT NOT NULL,
                    source_path TEXT NOT NULL,
                    metadata_json TEXT,
                    indexed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            """)
            conn.commit()

    def _load_or_index_knowledge(self):
        """Loads documents from database or indexes root workspace."""
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT doc_id, title, category, content, source_path, metadata_json FROM alfred_knowledge_docs")
            rows = cursor.fetchall()

        if rows:
            self.documents = [
                KnowledgeDocument(
                    doc_id=r[0], title=r[1], category=r[2],
                    content=r[3], source_path=r[4],
                    metadata=json.loads(r[5]) if r[5] else {}
                )
                for r in rows
            ]
            self._fit_vectorizer()
        else:
            self.index_workspace()

    def index_workspace(self, root_dir: str = ".") -> int:
        """Indexes academic lecture PPTs, DSA code, and system documentation."""
        docs: List[KnowledgeDocument] = []
        root = Path(root_dir)

        # 1. Index Lecture Presentations (PPTX)
        pptx_files = [
            ("computer_system_architecture.pptx", "COMPUTER_SYSTEM_ARCHITECTURE", "CSA Lecture Master Slide Deck"),
            ("Computer_System_Architecture_OS.pptx", "OPERATING_SYSTEMS", "OS & Memory Hierarchy Slide Deck"),
            ("Computer_System_Architecture_Cinematic.pptx", "COMPUTER_SYSTEM_ARCHITECTURE", "Cinematic Computer Architecture Concepts"),
        ]
        try:
            from pptx import Presentation
            for rel_path, cat, title in pptx_files:
                fp = root / rel_path
                if fp.exists():
                    try:
                        prs = Presentation(str(fp))
                        slide_texts = []
                        for s_idx, slide in enumerate(prs.slides):
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    slide_texts.append(shape.text_frame.text)
                        full_pptx_text = "\n".join(slide_texts)
                        if full_pptx_text.strip():
                            # Chunk by 1200 characters
                            chunks = [full_pptx_text[i:i+1200] for i in range(0, len(full_pptx_text), 1000)]
                            for idx, chunk in enumerate(chunks):
                                docs.append(KnowledgeDocument(
                                    doc_id=f"pptx_{rel_path.replace('.', '_')}_chunk_{idx}",
                                    title=f"{title} (Section {idx+1})",
                                    category=cat, content=chunk, source_path=str(fp),
                                    metadata={"type": "lecture_slide", "slide_count": len(prs.slides)}
                                ))
                    except Exception as e:
                        logger.warning(f"Failed to read pptx {fp}: {e}")
        except ImportError:
            pass

        # 2. Index DSA & Numpy Code
        code_targets = [
            ("dsa/day_1_arrays_and_hashmaps.py", "DATA_STRUCTURES", "DSA Day 1: Arrays, HashMaps & Two Pointers"),
            ("numpy_01_matrix_multiplication.py", "NUMPY_MATH", "Numpy 01: Matrix Multiplication & Linear Algebra"),
            ("numpy_02_upper_triangular.py", "NUMPY_MATH", "Numpy 02: Upper Triangular Matrices & Decompositions"),
            ("numpy_03_rank_det_trace_inv.py", "NUMPY_MATH", "Numpy 03: Rank, Determinant, Trace & Matrix Inverses"),
            ("numpy_04_slicing_joining_sorting.py", "NUMPY_MATH", "Numpy 04: Advanced Slicing, Joining & Sorting"),
            ("numpy_05_uniform_binomial.py", "NUMPY_MATH", "Numpy 05: Statistical Distributions (Uniform & Binomial)"),
            ("array_implementation.py", "DATA_STRUCTURES", "Custom Dynamic Array Implementation"),
        ]
        for rel_path, cat, title in code_targets:
            fp = root / rel_path
            if fp.exists():
                try:
                    text = fp.read_text(encoding="utf-8", errors="ignore")
                    docs.append(KnowledgeDocument(
                        doc_id=f"code_{rel_path.replace('/', '_').replace('.', '_')}",
                        title=title, category=cat, content=text, source_path=str(fp),
                        metadata={"type": "source_code"}
                    ))
                except Exception as e:
                    logger.warning(f"Failed to read code {fp}: {e}")

        # 3. Index Core System & Governance
        sys_files = [
            ("Alfred.Modelfile", "ALFRED_CORE", "Alfred Sovereign Butler & Engineer Modelfile"),
            ("src/jarvisx/cron/morning_routine.py", "SYSTEM_GOVERNANCE", "Autonomous Morning Wake-Up Routine"),
            ("src/jarvisx/executive/daily_executive.py", "ACADEMIC_10CGPA", "Daily Executive & 10-CGPA Schedule Sentinel"),
            ("src/jarvisx/developer/code_healer.py", "FRIDAY_DEV_CORE", "Friday Autonomous Code Healer & Test Synthesizer"),
        ]
        for rel_path, cat, title in sys_files:
            fp = root / rel_path
            if fp.exists():
                try:
                    text = fp.read_text(encoding="utf-8", errors="ignore")
                    docs.append(KnowledgeDocument(
                        doc_id=f"sys_{rel_path.replace('/', '_').replace('.', '_')}",
                        title=title, category=cat, content=text[:1500], source_path=str(fp),
                        metadata={"type": "system_module"}
                    ))
                except Exception as e:
                    logger.warning(f"Failed to read sys file {fp}: {e}")

        # 4. Store in SQLite Database
        with sqlite3.connect(self.db_path) as conn:
            conn.execute("DELETE FROM alfred_knowledge_docs")
            for doc in docs:
                conn.execute(
                    "INSERT INTO alfred_knowledge_docs (doc_id, title, category, content, source_path, metadata_json) VALUES (?, ?, ?, ?, ?, ?)",
                    (doc.doc_id, doc.title, doc.category, doc.content, doc.source_path, json.dumps(doc.metadata))
                )
            conn.commit()

        self.documents = docs
        self._fit_vectorizer()
        logger.info(f"Indexed {len(docs)} pure Alfred knowledge chunks into {self.db_path}")
        return len(docs)

    def _fit_vectorizer(self):
        """Fits TF-IDF vectorizer over all document contents."""
        if not self.documents:
            return
        corpus = [f"{d.title} {d.category} {d.content}" for d in self.documents]
        self.vectorizer = TfidfVectorizer(
            ngram_range=(1, 2),
            stop_words="english",
            max_features=10000
        )
        self.tfidf_matrix = self.vectorizer.fit_transform(corpus)

    def query(self, query_text: str, top_k: int = 3, min_score: float = 0.05) -> List[SearchResult]:
        """Performs fast semantic vector search over the knowledge base."""
        if not self.vectorizer or self.tfidf_matrix is None or not self.documents:
            return []

        query_vec = self.vectorizer.transform([query_text])
        similarities = cosine_similarity(query_vec, self.tfidf_matrix).flatten()

        top_indices = np.argsort(similarities)[::-1][:top_k]
        results = []
        for idx in top_indices:
            score = float(similarities[idx])
            if score < min_score:
                continue
            doc = self.documents[idx]
            snippet = doc.content[:300].replace("\n", " ").strip() + "..."
            results.append(SearchResult(
                doc_id=doc.doc_id,
                title=doc.title,
                category=doc.category,
                snippet=snippet,
                score=round(score, 4),
                source_path=doc.source_path,
                metadata=doc.metadata
            ))
        return results

    def get_stats(self) -> Dict[str, Any]:
        """Returns knowledge base statistics."""
        categories = {}
        for d in self.documents:
            categories[d.category] = categories.get(d.category, 0) + 1
        return {
            "total_indexed_chunks": len(self.documents),
            "db_path": self.db_path,
            "categories": categories,
            "vocabulary_size": len(self.vectorizer.vocabulary_) if self.vectorizer else 0
        }
