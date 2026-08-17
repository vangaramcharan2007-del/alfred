"""
scripts/generate_cinematic_ppt.py
Generates the 15-slide Luxury Minimalist Tech PowerPoint presentation with:
- Native OpenXML Morph & Smooth Transitions on all slides
- Cinematic Ambient Lighting & Glow Overlays
- Frosted Glass Cards with High-Tech Corner Framing Accents
- 3D Isometric Hardware Artworks (Motherboard, 3-Part Bus, Storage Pyramid, Multicore Chip)
- Structured 3-Speaker Content & Spoken Scripts:
    * Act I (Slides 1-5): V. Ram Charan — Core Hardware, CPU/RAM/IO, System Bus (Address, Data, Control), Interrupts & DMA
    * Act II (Slides 6-10): Vedhanth — Memory & Storage Hierarchy, Locality of Reference, Volatility, OS Optimizations
    * Act III (Slides 11-15): Lochan — Single vs Multiprocessor, SMP vs AMP, Multicore & Clustered Systems, OS Challenges
- Zero textbook or author citations
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

ROOT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS_DIR = os.path.join(ROOT_DIR, "assets", "presentation")
GAMMA_DIR = os.path.join(ASSETS_DIR, "gamma_assets")
sys.path.insert(0, ROOT_DIR)

from assets.presentation.gamma_illustrations import generate_all_gamma_assets

# Luxury Dark Palette (Gamma / Envato Style)
BG_CARBON = RGBColor(18, 20, 28)        # #12141C Deep Charcoal
CARD_BG = RGBColor(26, 30, 42)          # #1A1E2A Translucent Dark Card
CARD_BORDER = RGBColor(46, 54, 74)      # #2E364A Subtle Outline
CARD_HIGHLIGHT = RGBColor(38, 45, 64)   # #262D40 Highlight Card
TEXT_WHITE = RGBColor(255, 255, 255)    # #FFFFFF Crisp Pure White
TEXT_MUTED = RGBColor(160, 175, 200)    # #A0AFC8 Soft Slate
TEXT_DIM = RGBColor(110, 125, 150)      # #6E7D96 Dim Metadata
ACCENT_CYAN = RGBColor(56, 189, 248)    # #38BDF8 Soft Ice Cyan
ACCENT_INDIGO = RGBColor(129, 140, 248) # #818CF8 Hyper Indigo
ACCENT_EMERALD = RGBColor(52, 211, 153) # #34D399 Emerald Green
ACCENT_AMBER = RGBColor(245, 158, 11)   # #F59E0B Amber
ACCENT_ROSE = RGBColor(244, 63, 94)     # #F43F5E Coral / Alert
PILL_BG = RGBColor(30, 36, 52)          # #1E2434

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"

def set_slide_backdrop(slide, color=BG_CARBON):
    """Creates a clean solid matte carbon backdrop."""
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()

def apply_slide_morph_transition(slide):
    """Applies valid OpenXML Morph & Smooth Slide Transition."""
    tr_xml = parse_xml(
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" spd="med" dur="1500">'
        '<p14:morph option="byObject"/>'
        '</p:transition>'
    )
    slide._element.append(tr_xml)

def add_header(slide, slide_num, total_slides, act_title, speaker_name, title, subtitle=None):
    """Creates a clean Gamma-style header with subtle pill tags."""
    # Top Left Act Tag
    act_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.7), Inches(0.34))
    act_box.fill.solid()
    act_box.fill.fore_color.rgb = PILL_BG
    act_box.line.color.rgb = CARD_BORDER
    act_box.line.width = Pt(1)
    tf_act = act_box.text_frame
    tf_act.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_act = tf_act.paragraphs[0]
    p_act.text = act_title.upper()
    p_act.font.name = FONT_HEADING
    p_act.font.size = Pt(9)
    p_act.font.bold = True
    p_act.font.color.rgb = ACCENT_CYAN
    p_act.alignment = PP_ALIGN.CENTER

    # Top Right Speaker Badge
    spk_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(0.4), Inches(2.7), Inches(0.34))
    spk_box.fill.solid()
    spk_box.fill.fore_color.rgb = PILL_BG
    spk_box.line.color.rgb = CARD_BORDER
    spk_box.line.width = Pt(1)
    tf_spk = spk_box.text_frame
    tf_spk.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_spk = tf_spk.paragraphs[0]
    p_spk.text = speaker_name.upper()
    p_spk.font.name = FONT_HEADING
    p_spk.font.size = Pt(9)
    p_spk.font.bold = True
    p_spk.font.color.rgb = ACCENT_INDIGO
    p_spk.alignment = PP_ALIGN.CENTER

    # Top Right Counter
    num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.45), Inches(0.4), Inches(1.1), Inches(0.34))
    num_box.fill.solid()
    num_box.fill.fore_color.rgb = PILL_BG
    num_box.line.color.rgb = CARD_BORDER
    num_box.line.width = Pt(1)
    tf_num = num_box.text_frame
    tf_num.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_num = tf_num.paragraphs[0]
    p_num.text = f"{slide_num:02d} / {total_slides:02d}"
    p_num.font.name = FONT_HEADING
    p_num.font.size = Pt(9)
    p_num.font.bold = True
    p_num.font.color.rgb = TEXT_DIM
    p_num.alignment = PP_ALIGN.CENTER

    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.75))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_t = tf_title.paragraphs[0]
    p_t.text = title
    p_t.font.name = FONT_HEADING
    p_t.font.size = Pt(24)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE

    if subtitle:
        p_sub = tf_title.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(11.5)
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.space_before = Pt(2)

def add_card(slide, left, top, width, height, title=None, border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=None, corner_accent=True):
    """Creates a minimalist glass card with high-tech corner framing and glowing accents."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)

    if accent_bar:
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(height))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_bar
        bar.line.fill.background()

    if corner_accent:
        c_tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + width - 0.25), Inches(top + 0.08), Inches(0.15), Inches(0.02))
        c_tick.fill.solid()
        c_tick.fill.fore_color.rgb = accent_bar if accent_bar else CARD_BORDER
        c_tick.line.fill.background()

    if title:
        tb = slide.shapes.add_textbox(Inches(left + (0.24 if accent_bar else 0.2)), Inches(top + 0.15), Inches(width - 0.4), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_HEADING
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    return card

def add_bottom_banner(slide, text, tag="Key Principle:", tag_color=ACCENT_CYAN):
    """Adds a sleek full-width bottom callout banner matching the Gamma template."""
    banner = add_card(slide, 0.8, 6.35, 11.733, 0.65, border_color=CARD_BORDER, bg_color=PILL_BG)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    
    r_tag = p.add_run()
    r_tag.text = f"□ {tag} "
    r_tag.font.name = FONT_HEADING
    r_tag.font.size = Pt(10.5)
    r_tag.font.bold = True
    r_tag.font.color.rgb = tag_color
    
    r_body = p.add_run()
    r_body.text = text
    r_body.font.name = FONT_BODY
    r_body.font.size = Pt(10.5)
    r_body.font.color.rgb = TEXT_WHITE

def add_circular_step(slide, cx, cy, radius, label, sublabel=None, number=None, border_color=ACCENT_CYAN):
    """Draws sleek circular process step nodes matching Gamma template."""
    r_inch = Inches(radius)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - radius), Inches(cy - radius), Inches(radius * 2), Inches(radius * 2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PILL_BG
    circle.line.color.rgb = border_color
    circle.line.width = Pt(2)

    tb = slide.shapes.add_textbox(Inches(cx - radius - 0.2), Inches(cy - 0.2), Inches(radius * 2 + 0.4), Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if number:
        p.text = str(number)
        p.font.name = FONT_HEADING
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
    else:
        p.text = "< / >"
        p.font.name = FONT_HEADING
        p.font.size = Pt(12)
        p.font.color.rgb = ACCENT_CYAN

    if label:
        tb_lbl = slide.shapes.add_textbox(Inches(cx - 0.8), Inches(cy + radius + 0.1), Inches(1.6), Inches(0.6))
        tf_l = tb_lbl.text_frame
        tf_l.word_wrap = True
        p_l = tf_l.paragraphs[0]
        p_l.alignment = PP_ALIGN.CENTER
        p_l.text = label
        p_l.font.name = FONT_HEADING
        p_l.font.size = Pt(10)
        p_l.font.bold = True
        p_l.font.color.rgb = TEXT_WHITE

def set_speaker_notes(slide, what_to_say, concept, transition, cue):
    """Sets structured presenter script notes on the slide."""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = f"=== PRESENTER SCRIPT ===\n\n" \
                      f"1. WHAT TO SAY (Spoken Script):\n{what_to_say}\n\n" \
                      f"2. KEY ARCHITECTURAL CONCEPT:\n{concept}\n\n" \
                      f"3. TRANSITION TO NEXT SLIDE:\n{transition}\n\n" \
                      f"4. PRESENTATION CUE:\n{cue}"

def build_presentation():
    # 1. Generate all isometric assets with ambient lighting
    generate_all_gamma_assets()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: TITLE & PRESENTERS (V. Ram Charan)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s1)
    apply_slide_morph_transition(s1)

    tb_t = s1.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(7.0), Inches(2.2))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p1 = tf_t.paragraphs[0]
    p1.text = "Computer System Architecture"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE

    p2 = tf_t.add_paragraph()
    p2.text = "Core Hardware, Storage Hierarchies & Modern Multiprocessing"
    p2.font.name = FONT_HEADING
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_CYAN
    p2.space_before = Pt(6)

    p3 = tf_t.add_paragraph()
    p3.text = "Operating Systems — System Architecture & Execution Foundations"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(4)

    presenters = [
        ("V. Ram Charan", "Act I Lead — Core Hardware & System Bus", ACCENT_CYAN),
        ("Vedhanth", "Act II Lead — Memory & Storage Hierarchy", ACCENT_INDIGO),
        ("Lochan", "Act III Lead — Multiprocessing & Modern Architectures", ACCENT_EMERALD)
    ]
    for i, (name, role, col) in enumerate(presenters):
        py = 3.3 + i * 0.95
        avatar = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(py), Inches(0.7), Inches(0.7))
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = PILL_BG
        avatar.line.color.rgb = CARD_BORDER
        avatar.line.width = Pt(1)
        tf_av = avatar.text_frame
        p_av = tf_av.paragraphs[0]
        p_av.alignment = PP_ALIGN.CENTER
        p_av.text = "👤"
        p_av.font.size = Pt(14)
        
        tb_p = s1.shapes.add_textbox(Inches(1.65), Inches(py + 0.05), Inches(5.5), Inches(0.6))
        tf_p = tb_p.text_frame
        p_pn = tf_p.paragraphs[0]
        p_pn.text = name
        p_pn.font.name = FONT_HEADING
        p_pn.font.size = Pt(14)
        p_pn.font.bold = True
        p_pn.font.color.rgb = TEXT_WHITE
        p_pr = tf_p.add_paragraph()
        p_pr.text = role
        p_pr.font.name = FONT_BODY
        p_pr.font.size = Pt(10)
        p_pr.font.color.rgb = col

    q_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.3), Inches(0.06), Inches(0.65))
    q_bar.fill.solid()
    q_bar.fill.fore_color.rgb = ACCENT_CYAN
    q_bar.line.fill.background()
    tb_q = s1.shapes.add_textbox(Inches(0.95), Inches(6.3), Inches(6.8), Inches(0.65))
    tf_q = tb_q.text_frame
    tf_q.word_wrap = True
    pq = tf_q.paragraphs[0]
    pq.text = '"The operating system is the software most intimately involved with computer hardware."'
    pq.font.name = FONT_BODY
    pq.font.size = Pt(11)
    pq.font.italic = True
    pq.font.color.rgb = TEXT_MUTED

    mb_img = os.path.join(GAMMA_DIR, "iso_motherboard.png")
    if os.path.exists(mb_img):
        s1.shapes.add_picture(mb_img, Inches(7.5), Inches(0.8), width=Inches(5.2))

    set_speaker_notes(
        s1,
        "Good morning everyone. I’ll be explaining the core hardware of a computer system. Together with Vedhanth and Lochan, we will walk you through how hardware components communicate, how memory hierarchies optimize speed and capacity, and how modern multiprocessing systems scale concurrent execution.",
        "Computer system architecture establishes the hardware foundation for operating system process management, memory allocation, and I/O handling.",
        "Let us begin with Slide 2 by examining the three core hardware components: the CPU, main memory, and I/O devices.",
        "Stand center, introduce team members, gesture to the presenter cards."
    )

    # =========================================================================
    # SLIDE 2: CORE HARDWARE COMPONENTS (V. Ram Charan)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s2)
    apply_slide_morph_transition(s2)
    add_header(s2, 2, 15, "Act I: Core Hardware", "V. Ram Charan", "Core Hardware of a Computer System", "The Three Fundamental Hardware Pillars Powering Execution")

    col_w = 3.7
    hw_cards = [
        ("The CPU (Processor)", ACCENT_CYAN, [
            ("Instruction Cycle:", "Responsible for fetching, decoding, and executing instructions from memory."),
            ("Core Units:", "Contains the Arithmetic Logic Unit (ALU), Control Unit (CU), and internal registers."),
            ("System Driver:", "Directs computation and orchestrates data flow across all subsystems.")
        ]),
        ("Main Memory (RAM)", ACCENT_INDIGO, [
            ("Primary Storage:", "Stores data and active programs that are currently being used by the system."),
            ("Direct Access:", "The only large memory directly addressable and executable by the CPU."),
            ("Volatile Nature:", "Provides high-speed working storage for active runtime processes.")
        ]),
        ("I/O Devices & Controllers", ACCENT_EMERALD, [
            ("External Bridge:", "Helps the computer communicate with the outside world and users."),
            ("Device Categories:", "Displays, keyboards, network interfaces, and secondary storage."),
            ("Dedicated Hardware:", "Device controllers manage physical hardware protocols autonomously.")
        ])
    ]
    for i, (ctitle, ccol, citems) in enumerate(hw_cards):
        cx = 0.8 + i * 4.0
        add_card(s2, cx, 1.6, col_w, 4.4, ctitle, border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ccol)
        tb = s2.shapes.add_textbox(Inches(cx + 0.24), Inches(2.1), Inches(col_w - 0.45), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, (h, b) in enumerate(citems):
            p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
            p.text = "• " + h + " "
            p.font.name = FONT_HEADING
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            if j > 0: p.space_before = Pt(6)
            run = p.add_run()
            run.text = b
            run.font.name = FONT_BODY
            run.font.size = Pt(10.5)
            run.font.bold = False
            run.font.color.rgb = TEXT_MUTED

    add_bottom_banner(s2, "The CPU continuously executes the Fetch-Decode-Execute cycle on programs held in Main Memory, communicating results via I/O.", "Execution Foundation:", ACCENT_CYAN)

    set_speaker_notes(
        s2,
        "A computer system mainly consists of three parts: the CPU, main memory, and I/O devices. The CPU is responsible for fetching, decoding, and executing instructions. Main memory stores the data and programs that are currently being used, while I/O devices help the computer communicate with the outside world.",
        "The CPU, RAM, and I/O form the triad of computer architecture; the CPU cannot run code that is not first loaded into main memory.",
        "Now let's examine how these three components communicate with each other over the system bus.",
        "Point to the 3 distinct cards highlighting CPU, Main Memory, and I/O Devices."
    )

    # =========================================================================
    # SLIDE 3: THE SYSTEM BUS ARCHITECTURE (V. Ram Charan)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s3)
    apply_slide_morph_transition(s3)
    add_header(s3, 3, 15, "Act I: Core Hardware", "V. Ram Charan", "The System Bus Architecture", "Three Specialized High-Speed Communication Channels")

    # Left: 3 Dedicated Bus Cards
    bus_cards = [
        ("1. Address Bus", "Tells WHERE the data should go by carrying physical memory addresses and I/O port numbers (Unidirectional from CPU/DMA).", ACCENT_CYAN),
        ("2. Data Bus", "Carries the ACTUAL DATA and binary instructions between the CPU, main memory, and device controllers (Bidirectional).", ACCENT_INDIGO),
        ("3. Control Bus", "Carries CONTROL SIGNALS and timing clocks — such as Memory Read/Write, I/O Read/Write, and Bus Grant commands (Bidirectional).", ACCENT_EMERALD)
    ]
    for i, (btitle, bdesc, bcol) in enumerate(bus_cards):
        by = 1.6 + i * 1.5
        add_card(s3, 0.8, by, 5.8, 1.35, btitle, border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=bcol)
        tb = s3.shapes.add_textbox(Inches(1.05), Inches(by + 0.45), Inches(5.3), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bdesc
        p.font.name = FONT_BODY
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_MUTED

    # Right: 3-Part Bus Isometric Diagram
    bus_img = os.path.join(GAMMA_DIR, "iso_system_bus.png")
    if os.path.exists(bus_img):
        s3.shapes.add_picture(bus_img, Inches(7.0), Inches(1.5), width=Inches(5.5))

    add_bottom_banner(s3, "These three buses work synchronously to coordinate all instruction fetching, operand reads, and device data transfers.", "System Bus Synchronization:", ACCENT_CYAN)

    set_speaker_notes(
        s3,
        "These components communicate through a system bus, which has three parts: the address bus, which tells where the data should go; the data bus, which carries the actual data; and the control bus, which carries control signals.",
        "The system bus connects CPU, memory, and I/O controllers; bus width determines addressable memory limits and throughput bandwidth.",
        "Next, let's look at two critical concepts in hardware communication: Interrupts and DMA.",
        "Trace the Address, Data, and Control lines on the diagram."
    )

    # =========================================================================
    # SLIDE 4: INTERRUPTS & DMA (V. Ram Charan)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s4)
    apply_slide_morph_transition(s4)
    add_header(s4, 4, 15, "Act I: Core Hardware", "V. Ram Charan", "Hardware Communication: Interrupts & DMA", "Asynchronous Event Signaling & High-Speed Memory Transfers")

    # Left: 2 Hero Cards for Interrupts and DMA
    add_card(s4, 0.8, 1.6, 5.7, 4.4, "Interrupt Mechanism", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_CYAN)
    tb_int = s4.shapes.add_textbox(Inches(1.05), Inches(2.1), Inches(5.2), Inches(3.7))
    tf_int = tb_int.text_frame
    tf_int.word_wrap = True
    pts_int = [
        ("Purpose:", "Allows devices to get the CPU’s attention asynchronously whenever they need it."),
        ("Eliminates Polling:", "CPU avoids wasting billions of clock cycles in busy-waiting polling loops."),
        ("Vector Table Dispatch:", "CPU pushes state to stack and branches instantly to the Interrupt Service Routine (ISR)."),
        ("Seamless Resumption:", "User execution resumes immediately after device service completes.")
    ]
    for i, (h, b) in enumerate(pts_int):
        p = tf_int.add_paragraph() if i > 0 else tf_int.paragraphs[0]
        p.text = "• " + h + " "
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        if i > 0: p.space_before = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_card(s4, 6.8, 1.6, 5.7, 4.4, "Direct Memory Access (DMA)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_EMERALD)
    tb_dma = s4.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(5.2), Inches(3.7))
    tf_dma = tb_dma.text_frame
    tf_dma.word_wrap = True
    pts_dma = [
        ("Purpose:", "Allows devices to transfer data directly to memory without making the CPU handle every piece of data."),
        ("High-Speed Channels:", "Essential for disk controllers, SSDs, graphics, and Gigabit network interfaces."),
        ("Single Interrupt per Block:", "Transfers entire data blocks directly between device buffer and RAM; emits only 1 interrupt per block."),
        ("CPU Offloading:", "Frees the CPU to execute user computation while memory transfers proceed in parallel.")
    ]
    for i, (h, b) in enumerate(pts_dma):
        p = tf_dma.add_paragraph() if i > 0 else tf_dma.paragraphs[0]
        p.text = "• " + h + " "
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        if i > 0: p.space_before = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_bottom_banner(s4, "Interrupts provide asynchronous control signaling, while DMA provides high-throughput direct data movement.", "Core Takeaway:", ACCENT_EMERALD)

    set_speaker_notes(
        s4,
        "Two important concepts are interrupts and DMA. Interrupts allow devices to get the CPU’s attention when they need it. DMA allows devices to transfer data directly to memory without making the CPU handle every piece of data.",
        "Interrupts eliminate busy-wait polling overhead; DMA offloads byte-by-byte data transfer from the CPU.",
        "Let us summarize how core hardware components integrate before moving to memory hierarchy.",
        "Contrast the event-driven signal flow of interrupts with the high-throughput block stream of DMA."
    )

    # =========================================================================
    # SLIDE 5: CORE HARDWARE INTEGRATION SUMMARY (V. Ram Charan)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s5)
    apply_slide_morph_transition(s5)
    add_header(s5, 5, 15, "Act I: Core Hardware", "V. Ram Charan", "Core Hardware Integration Summary", "Cohesive Operation Across Processor, Memory, Bus & Controllers")

    # 3 Summary Cards
    col_w = 3.7
    int_cards = [
        ("1. Triad Organization", ACCENT_CYAN, "CPU fetches, decodes, and executes. Main memory holds active tasks. I/O bridges external devices into the system."),
        ("2. Tri-Bus Highway", ACCENT_INDIGO, "Address bus routes destinations, data bus transmits operands/code, and control bus synchronizes read/write timing."),
        ("3. Event & Data Flow", ACCENT_EMERALD, "Interrupts notify the CPU on demand, while DMA streams high-speed data directly to RAM without processor bottlenecks.")
    ]
    for i, (stitle, scol, sdesc) in enumerate(int_cards):
        cx = 0.8 + i * 4.0
        add_card(s5, cx, 1.6, col_w, 4.4, stitle, border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=scol)
        tb = s5.shapes.add_textbox(Inches(cx + 0.24), Inches(2.2), Inches(col_w - 0.45), Inches(3.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sdesc
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(4)

    add_bottom_banner(s5, "That’s the basic idea of how the hardware components communicate. Next: Memory & Storage Hierarchy.", "Act I Transition:", ACCENT_CYAN)

    set_speaker_notes(
        s5,
        "That’s the basic idea of how the hardware components communicate. I will now pass the presentation to Vedhanth, who will explain the memory and storage hierarchy, locality of reference, and how the operating system manages storage tiers efficiently.",
        "Hardware components form a tightly synchronized system coordinated by the system bus, interrupts, and DMA.",
        "Pass presentation to Vedhanth for Act II.",
        "Hand over clicker to Vedhanth."
    )

    # =========================================================================
    # SLIDE 6: MEMORY & STORAGE HIERARCHY OVERVIEW (Vedhanth)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s6)
    apply_slide_morph_transition(s6)
    add_header(s6, 6, 15, "Act II: Memory & Storage", "Vedhanth", "Memory and Storage Hierarchy", "Balancing Speed, Cost, and Capacity Across Storage Tiers")

    # Left: Explanation & Principle
    tb_mh = s6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.5))
    tf_mh = tb_mh.text_frame
    tf_mh.word_wrap = True

    p_mh1 = tf_mh.paragraphs[0]
    p_mh1.text = "Why Do Computers Use Different Types of Memory?"
    p_mh1.font.name = FONT_HEADING
    p_mh1.font.size = Pt(15)
    p_mh1.font.bold = True
    p_mh1.font.color.rgb = TEXT_WHITE

    p_mh2 = tf_mh.add_paragraph()
    p_mh2.text = "Computers use different types of memory because faster memory is usually smaller and more expensive, while slower memory can provide much larger capacity."
    p_mh2.font.name = FONT_BODY
    p_mh2.font.size = Pt(11)
    p_mh2.font.color.rgb = ACCENT_CYAN
    p_mh2.space_before = Pt(6)

    # 2 Sub Cards
    add_card(s6, 0.8, 3.1, 5.8, 1.4, "High-Speed Tier (Top)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_CYAN)
    tb_t1 = s6.shapes.add_textbox(Inches(1.05), Inches(3.55), Inches(5.3), Inches(0.85))
    p = tb_t1.text_frame.paragraphs[0]
    p.text = "Registers and CPU caches provide sub-nanosecond access speeds to keep the execution pipeline fed, but are limited in physical capacity."
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED

    add_card(s6, 0.8, 4.65, 5.8, 1.4, "High-Capacity Tier (Bottom)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_AMBER)
    tb_t2 = s6.shapes.add_textbox(Inches(1.05), Inches(5.1), Inches(5.3), Inches(0.85))
    p = tb_t2.text_frame.paragraphs[0]
    p.text = "Main memory and secondary storage (SSDs/HDDs) provide gigabytes to terabytes of storage at low cost per bit."
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED

    # Right: Storage Pyramid Illustration
    pyr_img = os.path.join(GAMMA_DIR, "iso_storage_pyramid.png")
    if os.path.exists(pyr_img):
        s6.shapes.add_picture(pyr_img, Inches(7.0), Inches(1.5), width=Inches(5.5))

    add_bottom_banner(s6, "Faster memory is smaller and expensive; slower storage is large and economical. The hierarchy balances both.", "Fundamental Trade-off:", ACCENT_CYAN)

    set_speaker_notes(
        s6,
        "Now I’ll explain memory and storage hierarchy. Computers use different types of memory because faster memory is usually smaller and more expensive, while slower memory can provide much larger capacity. The hierarchy allows the computer to approximate the speed of the fastest memory with the capacity of the largest storage.",
        "Trade-off between access latency, cost per bit, and storage capacity across the memory hierarchy.",
        "Let us examine each tier of the hierarchy from internal registers down to secondary storage.",
        "Trace the pyramid from top (Registers/Cache) to bottom (RAM/SSD)."
    )

    # =========================================================================
    # SLIDE 7: HIERARCHY TIERS (Vedhanth)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s7)
    apply_slide_morph_transition(s7)
    add_header(s7, 7, 15, "Act II: Memory & Storage", "Vedhanth", "The Hierarchy Tiers: Registers to Secondary Storage", "Speed, Capacity, and Access Latency Across the Four Levels")

    c_w4 = 2.75
    tiers_data = [
        ("1. Registers", ACCENT_CYAN, "< 1 ns", "Bytes / KB", "Fastest & smallest storage; located directly inside the CPU; holds immediate operands."),
        ("2. Cache (L1/L2/L3)", ACCENT_INDIGO, "1 – 10 ns", "KB – MB", "High-speed on-chip SRAM; buffers frequently used instructions and active data."),
        ("3. Main Memory (RAM)", ACCENT_EMERALD, "50 – 100 ns", "GBs", "Primary DRAM workspace; holds currently executing programs and operating system buffers."),
        ("4. Secondary Storage", ACCENT_AMBER, "10 μs – 10 ms", "TBs", "SSDs, NVMe & Hard Disks; non-volatile long-term storage preserving files permanently.")
    ]
    for i, (ttitle, tcol, tspeed, tcap, tdesc) in enumerate(tiers_data):
        tx = 0.8 + i * 3.0
        add_card(s7, tx, 1.6, c_w4, 4.4, ttitle, border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=tcol)
        
        # Speed & Cap Badges
        tb_b = s7.shapes.add_textbox(Inches(tx + 0.2), Inches(2.1), Inches(c_w4 - 0.4), Inches(0.8))
        tf_b = tb_b.text_frame
        p_s = tf_b.paragraphs[0]
        p_s.text = f"Speed: {tspeed}"
        p_s.font.name = FONT_HEADING
        p_s.font.size = Pt(10)
        p_s.font.bold = True
        p_s.font.color.rgb = tcol
        p_c = tf_b.add_paragraph()
        p_c.text = f"Capacity: {tcap}"
        p_c.font.name = FONT_BODY
        p_c.font.size = Pt(9.5)
        p_c.font.color.rgb = TEXT_WHITE

        tb_d = s7.shapes.add_textbox(Inches(tx + 0.2), Inches(3.0), Inches(c_w4 - 0.4), Inches(2.8))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = tdesc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED

    add_bottom_banner(s7, "The hierarchy starts with registers, which are the fastest and smallest, followed by cache, RAM, and finally secondary storage such as SSDs and hard disks.", "Hierarchy Structure:", ACCENT_CYAN)

    set_speaker_notes(
        s7,
        "The hierarchy starts with registers, which are the fastest and smallest, followed by cache, RAM, and finally secondary storage such as SSDs and hard disks. Registers operate at processor clock speeds, cache buffers hot data, RAM holds active programs, and secondary storage provides permanent capacity.",
        "Step-by-step characteristics of Registers, Cache (SRAM), Main Memory (DRAM), and Secondary Storage (NVM/Disks).",
        "Next, why does caching work so effectively? Let's look at the Locality of Reference.",
        "Point to the 4 columns from fastest (Registers) to largest (Secondary Storage)."
    )

    # =========================================================================
    # SLIDE 8: LOCALITY OF REFERENCE (Vedhanth)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s8)
    apply_slide_morph_transition(s8)
    add_header(s8, 8, 15, "Act II: Memory & Storage", "Vedhanth", "Locality of Reference: Temporal & Spatial", "The Behavioral Principle Making Caching and Hierarchy Work")

    # 2 Hero Cards: Temporal Locality vs Spatial Locality
    c_w = 5.7
    add_card(s8, 0.8, 1.6, c_w, 4.4, "Temporal Locality (Locality in Time)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_CYAN)
    tb_tem = s8.shapes.add_textbox(Inches(1.05), Inches(2.1), Inches(c_w - 0.4), Inches(3.7))
    tf_tem = tb_tem.text_frame
    tf_tem.word_wrap = True
    pts_tem = [
        ("Definition:", "Data used recently is likely to be used again soon in the near future."),
        ("Program Patterns:", "Loop counters, repeated function calls, stack variables, and active subroutines."),
        ("Cache Impact:", "Keeping recently accessed words in fast cache SRAM delivers hit rates above 90%."),
        ("Principle:", "If you read address X at time T, you will likely read address X again at time T + Δ.")
    ]
    for i, (h, b) in enumerate(pts_tem):
        p = tf_tem.add_paragraph() if i > 0 else tf_tem.paragraphs[0]
        p.text = "• " + h + " "
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        if i > 0: p.space_before = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_card(s8, 6.8, 1.6, c_w, 4.4, "Spatial Locality (Locality in Space)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_INDIGO)
    tb_spa = s8.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(c_w - 0.4), Inches(3.7))
    tf_spa = tb_spa.text_frame
    tf_spa.word_wrap = True
    pts_spa = [
        ("Definition:", "Data near recently used data is likely to be needed next by the program."),
        ("Program Patterns:", "Sequential instruction execution, traversing continuous arrays, and structured records."),
        ("Cache Impact:", "Memory controllers fetch whole cache lines (e.g. 64-byte blocks) rather than single words."),
        ("Principle:", "If you read address X, you will likely read address X + 1, X + 2 soon after.")
    ]
    for i, (h, b) in enumerate(pts_spa):
        p = tf_spa.add_paragraph() if i > 0 else tf_spa.paragraphs[0]
        p.text = "• " + h + " "
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        if i > 0: p.space_before = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_bottom_banner(s8, "Cache is important because programs show locality of reference: temporal (reuse soon) and spatial (nearby data needed next).", "Locality Rule:", ACCENT_CYAN)

    set_speaker_notes(
        s8,
        "Cache is important because programs usually show locality of reference. Temporal locality means that data used recently is likely to be used again soon. Spatial locality means that data near recently used data is likely to be needed next. Because of locality, small fast caches can satisfy the vast majority of memory accesses.",
        "Temporal locality (reuse in time) and Spatial locality (adjacency in space) are the core reasons caching succeeds.",
        "Now let's examine the difference between volatile and non-volatile storage.",
        "Highlight the contrast between loop re-execution (Temporal) and sequential array scans (Spatial)."
    )

    # =========================================================================
    # SLIDE 9: VOLATILITY VS PERMANENCE (Vedhanth)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s9)
    apply_slide_morph_transition(s9)
    add_header(s9, 9, 15, "Act II: Memory & Storage", "Vedhanth", "Memory Volatility: RAM vs. SSDs & Hard Disks", "Volatile Working Storage vs. Non-Volatile Permanent Retention")

    c_w = 5.7
    add_card(s9, 0.8, 1.6, c_w, 4.4, "Volatile Memory (RAM & Caches)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_ROSE)
    tb_vol = s9.shapes.add_textbox(Inches(1.05), Inches(2.1), Inches(c_w - 0.4), Inches(3.7))
    tf_vol = tb_vol.text_frame
    tf_vol.word_wrap = True
    pts_vol = [
        ("Core Characteristic:", "RAM is volatile, meaning it loses its contents when power is turned off."),
        ("Physical Storage:", "Uses dynamic capacitor charges (DRAM) or transistor latches (SRAM) requiring active power."),
        ("System Role:", "Acts as high-speed temporary workspace for currently running processes and kernel state."),
        ("Speed Advantage:", "Direct electrical access delivers nanosecond-level response times.")
    ]
    for i, (h, b) in enumerate(pts_vol):
        p = tf_vol.add_paragraph() if i > 0 else tf_vol.paragraphs[0]
        p.text = "• " + h + " "
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        if i > 0: p.space_before = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_card(s9, 6.8, 1.6, c_w, 4.4, "Non-Volatile Storage (SSDs & HDDs)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_EMERALD)
    tb_nvol = s9.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(c_w - 0.4), Inches(3.7))
    tf_nvol = tb_nvol.text_frame
    tf_nvol.word_wrap = True
    pts_nvol = [
        ("Core Characteristic:", "SSDs and hard disks are non-volatile, so they retain data even without power."),
        ("Physical Storage:", "Uses NAND flash memory floating gates (SSDs/NVMe) or magnetic platter domains (HDDs)."),
        ("System Role:", "Preserves operating system code, application files, and user data permanently."),
        ("Capacity Advantage:", "Provides massive terabyte-scale capacity at low cost per gigabyte.")
    ]
    for i, (h, b) in enumerate(pts_nvol):
        p = tf_nvol.add_paragraph() if i > 0 else tf_nvol.paragraphs[0]
        p.text = "• " + h + " "
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        if i > 0: p.space_before = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_bottom_banner(s9, "RAM is volatile (loses data on power-off); SSDs and HDDs are non-volatile (retain data permanently without power).", "Volatility Distinction:", ACCENT_ROSE)

    set_speaker_notes(
        s9,
        "RAM is volatile, meaning it loses its contents when power is turned off. SSDs and hard disks are non-volatile, so they retain data even without power. Because RAM is volatile, all programs and data must be saved to non-volatile secondary storage to survive reboots.",
        "Volatile working RAM vs Non-volatile secondary storage (SSDs/HDDs) and the role of permanent file storage.",
        "Next, how does the Operating System manage this entire storage hierarchy efficiently?",
        "Emphasize that volatile RAM provides execution speed while non-volatile SSDs provide data permanence."
    )

    # =========================================================================
    # SLIDE 10: OS STORAGE MANAGEMENT TECHNIQUES (Vedhanth)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s10)
    apply_slide_morph_transition(s10)
    add_header(s10, 10, 15, "Act II: Memory & Storage", "Vedhanth", "OS Memory Management Techniques", "Software Strategies for Maximizing Storage Hierarchy Performance")

    c_w = 5.7
    os_techs = [
        ("1. Paging & Virtual Memory", ACCENT_CYAN, "Divides memory into fixed-size pages and maps virtual addresses to physical RAM frames, swapping inactive pages to disk when memory is full."),
        ("2. Buffering", ACCENT_INDIGO, "Temporarily holds data in memory during device transfers to bridge speed mismatches between fast CPUs and slower I/O peripherals."),
        ("3. Disk Caching", ACCENT_EMERALD, "Maintains copies of frequently accessed file system blocks in RAM to avoid slow physical disk read operations."),
        ("4. Prefetching", ACCENT_AMBER, "Anticipates future memory and file requests based on spatial locality, loading data into cache/RAM before it is explicitly requested.")
    ]
    for i, (otitle, ocol, odesc) in enumerate(os_techs):
        gx = 0.8 if i % 2 == 0 else 6.8
        gy = 1.6 if i < 2 else 3.8
        add_card(s10, gx, gy, c_w, 2.0, otitle, border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ocol)
        tb = s10.shapes.add_textbox(Inches(gx + 0.24), Inches(gy + 0.5), Inches(c_w - 0.45), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = odesc
        p.font.name = FONT_BODY
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_MUTED

    add_bottom_banner(s10, "The OS uses techniques such as paging, buffering, disk caching, and prefetching to use this hierarchy efficiently.", "OS Optimization Suite:", ACCENT_CYAN)

    set_speaker_notes(
        s10,
        "The OS uses techniques such as paging, buffering, disk caching, and prefetching to use this hierarchy efficiently. Paging manages virtual address spaces, buffering smooths I/O data streams, disk caching keeps hot blocks in memory, and prefetching loads data ahead of time. That concludes Act II. I now pass to Lochan.",
        "Operating system storage management techniques: Paging, Buffering, Disk Caching, and Prefetching.",
        "I will now pass the presentation to Lochan, who will explain multiprocessing, modern multicore systems, and OS challenges.",
        "Hand over clicker to Lochan."
    )

    # =========================================================================
    # SLIDE 11: SINGLE-PROCESSOR VS MULTIPROCESSOR (Lochan)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s11)
    apply_slide_morph_transition(s11)
    add_header(s11, 11, 15, "Act III: Multiprocessing", "Lochan", "Single-Processor vs. Multiprocessor Systems", "Architectural Scaling from Single Execution Streams to Parallel Computing")

    # Left: Text descriptions
    tb_s11 = s11.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.5))
    tf_s11 = tb_s11.text_frame
    tf_s11.word_wrap = True

    p_sp = tf_s11.paragraphs[0]
    p_sp.text = "Single-Processor Systems"
    p_sp.font.name = FONT_HEADING
    p_sp.font.size = Pt(16)
    p_sp.font.bold = True
    p_sp.font.color.rgb = TEXT_WHITE

    pts_sp = [
        "Has one main CPU executing general-purpose instructions.",
        "All instructions must pass through that single CPU.",
        "Concurrency is simulated via rapid time-sliced context switching."
    ]
    for pt in pts_sp:
        p = tf_s11.add_paragraph()
        p.text = "• " + pt
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(3)

    p_mp = tf_s11.add_paragraph()
    p_mp.text = "Multiprocessor Systems"
    p_mp.font.name = FONT_HEADING
    p_mp.font.size = Pt(16)
    p_mp.font.bold = True
    p_mp.font.color.rgb = TEXT_WHITE
    p_mp.space_before = Pt(16)

    pts_mp = [
        "Has multiple CPUs sharing bus, clock, memory, and peripherals.",
        "Allows the OS to run more work at the same time (true hardware parallelism).",
        "Delivers higher system throughput, economy of scale, and graceful degradation."
    ]
    for pt in pts_mp:
        p = tf_s11.add_paragraph()
        p.text = "• " + pt
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(3)

    # Right: High-Contrast Tree Card
    tree_img = os.path.join(GAMMA_DIR, "single_vs_multi_tree.png")
    if os.path.exists(tree_img):
        s11.shapes.add_picture(tree_img, Inches(7.0), Inches(1.5), width=Inches(5.5))

    add_bottom_banner(s11, "A single-processor has one CPU for all instructions; a multiprocessor has multiple CPUs, running more work concurrently.", "Core Difference:", ACCENT_CYAN)

    set_speaker_notes(
        s11,
        "Finally, I’ll explain multiprocessing and modern computer systems. A single-processor system has one CPU, so all instructions must pass through that CPU. A multiprocessor system has multiple CPUs, allowing the OS to run more work at the same time.",
        "Single-processor serial execution vs multiprocessor concurrent hardware execution.",
        "Let us examine the two primary multiprocessing architectural approaches: SMP and AMP.",
        "Point to the single CPU path versus the multiple parallel CPU paths on the diagram."
    )

    # =========================================================================
    # SLIDE 12: MULTIPROCESSING: SMP VS AMP (Lochan)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s12)
    apply_slide_morph_transition(s12)
    add_header(s12, 12, 15, "Act III: Multiprocessing", "Lochan", "Multiprocessing Approaches: SMP vs. AMP", "Peer Multiprocessing vs. Master-Slave Hierarchical Scheduling")

    c_w = 5.7
    add_card(s12, 0.8, 1.6, c_w, 4.4, "Symmetric Multiprocessing (SMP)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_CYAN)
    tb_smp = s12.shapes.add_textbox(Inches(1.05), Inches(2.1), Inches(c_w - 0.4), Inches(3.7))
    tf_smp = tb_smp.text_frame
    tf_smp.word_wrap = True
    pts_smp = [
        ("Peer Relationship:", "All processors are treated as equal peers; no single processor is master."),
        ("Independent Execution:", "Each processor can run both user tasks and operating system kernel code concurrently."),
        ("Shared Resources:", "All CPUs share system memory and I/O channels over a common bus."),
        ("Modern Standard:", "Dominates modern computing from laptops and smartphones to large multi-socket servers.")
    ]
    for i, (h, b) in enumerate(pts_smp):
        p = tf_smp.add_paragraph() if i > 0 else tf_smp.paragraphs[0]
        p.text = "• " + h + " "
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        if i > 0: p.space_before = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_card(s12, 6.8, 1.6, c_w, 4.4, "Asymmetric Multiprocessing (AMP)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_AMBER)
    tb_amp = s12.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(c_w - 0.4), Inches(3.7))
    tf_amp = tb_amp.text_frame
    tf_amp.word_wrap = True
    pts_amp = [
        ("Master-Slave Model:", "One main master processor controls the system and assigns work to the others."),
        ("Centralized Control:", "Master processor runs the OS kernel and schedules tasks; slave processors execute assigned worker code."),
        ("Specialized Roles:", "Slaves may handle dedicated tasks like graphics, audio decoding, or real-time signal processing."),
        ("Simpler Scheduling:", "Avoids complex kernel data structure locking since only the master manages scheduling.")
    ]
    for i, (h, b) in enumerate(pts_amp):
        p = tf_amp.add_paragraph() if i > 0 else tf_amp.paragraphs[0]
        p.text = "• " + h + " "
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        if i > 0: p.space_before = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_bottom_banner(s12, "In SMP, all processors are peers and can run tasks independently; in AMP, a master processor assigns work to slaves.", "SMP vs AMP Distinction:", ACCENT_CYAN)

    set_speaker_notes(
        s12,
        "There are two approaches: SMP and AMP. In Symmetric Multiprocessing, or SMP, all processors are treated as peers and can run different tasks. In Asymmetric Multiprocessing, or AMP, one processor acts as the main processor and assigns work to the others.",
        "Symmetric (peer-to-peer) vs Asymmetric (master-slave) multiprocessing models and their trade-offs.",
        "Now let's examine modern systems: Multicore chips and Clustered systems.",
        "Contrast the decentralized peer model of SMP with the centralized master-slave model of AMP."
    )

    # =========================================================================
    # SLIDE 13: MULTICORE & CLUSTERED SYSTEMS (Lochan)
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s13)
    apply_slide_morph_transition(s13)
    add_header(s13, 13, 15, "Act III: Multiprocessing", "Lochan", "Modern Systems: Multicore & Clustered Systems", "On-Chip Integration (Scale-Up) vs. Networked Independent Machines (Scale-Out)")

    # Left: Multicore & Clustered Cards
    add_card(s13, 0.8, 1.6, 5.8, 2.1, "Multicore Processors (Scale-Up)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_CYAN)
    tb_mc = s13.shapes.add_textbox(Inches(1.05), Inches(2.05), Inches(5.3), Inches(1.5))
    tf_mc = tb_mc.text_frame
    tf_mc.word_wrap = True
    pts_mc = [
        "Multiple CPU cores are placed on a single physical chip.",
        "On-chip communication between cores is vastly faster than off-chip motherboard buses.",
        "Consumes significantly less power while sharing on-die L2/L3 caches."
    ]
    for pt in pts_mc:
        p = tf_mc.add_paragraph()
        p.text = "• " + pt
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(2)

    add_card(s13, 0.8, 3.9, 5.8, 2.1, "Clustered Systems (Scale-Out)", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_INDIGO)
    tb_cl = s13.shapes.add_textbox(Inches(1.05), Inches(4.35), Inches(5.3), Inches(1.5))
    tf_cl = tb_cl.text_frame
    tf_cl.word_wrap = True
    pts_cl = [
        "Multiple independent computer systems work together over a high-speed network.",
        "Nodes share a common Storage Area Network (SAN) for unified data access.",
        "Provides high availability and fault tolerance — if one machine fails, others take over."
    ]
    for pt in pts_cl:
        p = tf_cl.add_paragraph()
        p.text = "• " + pt
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(2)

    # Right: Multicore Chip 3D Art
    chip_img = os.path.join(GAMMA_DIR, "iso_multicore_chip.png")
    if os.path.exists(chip_img):
        s13.shapes.add_picture(chip_img, Inches(7.0), Inches(1.5), width=Inches(5.5))

    add_bottom_banner(s13, "Modern systems use multicore processors on single chips and clustered systems connecting independent computers over networks.", "Modern Architectural Spectrum:", ACCENT_CYAN)

    set_speaker_notes(
        s13,
        "Modern systems can also use multicore processors, where multiple CPU cores are placed on one chip. Another approach is clustered systems, where multiple independent computers work together over a network. Multicore scales up on a single die, while clustering scales out across many networked nodes.",
        "Multicore processors (on-chip parallelism) and Clustered systems (distributed high-availability nodes).",
        "However, multiprocessing creates significant challenges for the operating system. Let's look at those challenges.",
        "Contrast on-chip multicore scaling with distributed networked cluster scaling."
    )

    # =========================================================================
    # SLIDE 14: OS CHALLENGES IN MULTIPROCESSING (Lochan)
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s14)
    apply_slide_morph_transition(s14)
    add_header(s14, 14, 15, "Act III: Multiprocessing", "Lochan", "OS Challenges in Multiprocessing", "Critical Complexities in Coordinating Concurrent Hardware Execution")

    c_w = 5.7
    challenges = [
        ("1. CPU Scheduling", ACCENT_CYAN, "Deciding which processor or core runs which process to maximize throughput, avoid starvation, and maintain thread affinity."),
        ("2. Cache Coherence", ACCENT_INDIGO, "Ensuring that when one CPU modifies data in its private cache, other CPUs with copies of the same memory location are kept updated."),
        ("3. Load Balancing", ACCENT_EMERALD, "Distributing workload evenly across all available processors so no single CPU is bottlenecked while others sit idle."),
        ("4. Synchronization", ACCENT_ROSE, "Using locks, mutexes, and semaphores to protect shared kernel data structures and prevent race conditions or data corruption.")
    ]
    for i, (ctitle, ccol, cdesc) in enumerate(challenges):
        gx = 0.8 if i % 2 == 0 else 6.8
        gy = 1.6 if i < 2 else 3.8
        add_card(s14, gx, gy, c_w, 2.0, ctitle, border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ccol)
        tb = s14.shapes.add_textbox(Inches(gx + 0.24), Inches(gy + 0.5), Inches(c_w - 0.45), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cdesc
        p.font.name = FONT_BODY
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_MUTED

    add_bottom_banner(s14, "Multiprocessing creates challenges for the OS, including scheduling, cache coherence, load balancing, and synchronization.", "Operating System Requirement:", ACCENT_ROSE)

    set_speaker_notes(
        s14,
        "However, multiprocessing creates challenges for the OS, including scheduling, cache coherence, load balancing, and synchronization. So, multiprocessing allows computers to handle more work efficiently while requiring the OS to carefully coordinate everything.",
        "Four major OS multiprocessing challenges: Scheduling, Cache Coherence, Load Balancing, and Synchronization.",
        "Now let's bring our entire presentation together into our final synthesis and key takeaways.",
        "Review each of the four challenges and explain why hardware parallelism demands sophisticated OS coordination."
    )

    # =========================================================================
    # SLIDE 15: MASTER SYNTHESIS & KEY TAKEAWAYS (All 3 Presenters)
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_backdrop(s15)
    apply_slide_morph_transition(s15)
    add_header(s15, 15, 15, "Synthesis & Conclusion", "V. Ram Charan · Vedhanth · Lochan", "Complete Architecture & Key Takeaways", "Hardware Foundations, Storage Hierarchies & Modern Multiprocessing")

    # 5 Capsule Cards in a clean 4-top + 1-bottom grid
    capsules_top = [
        (1, "Core Hardware", "CPU (fetch/decode/exec), RAM, and I/O linked via Address, Data, and Control buses with Interrupts & DMA.", ACCENT_CYAN),
        (2, "Storage Hierarchy", "Balances speed, cost, and capacity from fast volatile Registers/Cache/RAM down to non-volatile SSDs/HDDs.", ACCENT_INDIGO),
        (3, "Locality of Reference", "Temporal (reuse soon) and Spatial (nearby) locality leveraged by Paging, Buffering, Caching & Prefetching.", ACCENT_EMERALD),
        (4, "Multiprocessing", "SMP (peers) and AMP (master-slave) scale hardware capacity alongside Multicore and Clustered systems.", ACCENT_AMBER)
    ]
    c_w4 = 2.75
    for num, ctitle, cdesc, ccol in capsules_top:
        cx = 0.8 + (num - 1) * 3.0
        add_card(s15, cx, 1.6, c_w4, 2.7, ctitle, border_color=CARD_BORDER, bg_color=CARD_BG)
        
        num_badge = s15.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + c_w4/2 - 0.25), Inches(1.4), Inches(0.5), Inches(0.5))
        num_badge.fill.solid()
        num_badge.fill.fore_color.rgb = PILL_BG
        num_badge.line.color.rgb = ccol
        num_badge.line.width = Pt(1.5)
        p_nb = num_badge.text_frame.paragraphs[0]
        p_nb.alignment = PP_ALIGN.CENTER
        p_nb.text = str(num)
        p_nb.font.name = FONT_HEADING
        p_nb.font.size = Pt(12)
        p_nb.font.bold = True
        p_nb.font.color.rgb = TEXT_WHITE

        tb = s15.shapes.add_textbox(Inches(cx + 0.15), Inches(2.0), Inches(c_w4 - 0.3), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cdesc
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED

    # Middle Row: 5th Wide Capsule Card (OS Multiprocessing Coordination)
    add_card(s15, 0.8, 4.6, 11.733, 1.0, "OS Multiprocessing Coordination", border_color=CARD_BORDER, bg_color=CARD_BG, accent_bar=ACCENT_ROSE)
    num_badge5 = s15.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8 + 11.733/2 - 0.25), Inches(4.4), Inches(0.5), Inches(0.5))
    num_badge5.fill.solid()
    num_badge5.fill.fore_color.rgb = PILL_BG
    num_badge5.line.color.rgb = ACCENT_ROSE
    num_badge5.line.width = Pt(1.5)
    p_nb5 = num_badge5.text_frame.paragraphs[0]
    p_nb5.alignment = PP_ALIGN.CENTER
    p_nb5.text = "5"
    p_nb5.font.name = FONT_HEADING
    p_nb5.font.size = Pt(12)
    p_nb5.font.bold = True
    p_nb5.font.color.rgb = TEXT_WHITE

    tb_5 = s15.shapes.add_textbox(Inches(1.1), Inches(4.9), Inches(11.2), Inches(0.6))
    tf_5 = tb_5.text_frame
    tf_5.word_wrap = True
    p5 = tf_5.paragraphs[0]
    p5.text = "Multiprocessing allows computers to handle more work efficiently while requiring the OS to carefully coordinate scheduling, cache coherence, load balancing, and synchronization."
    p5.font.name = FONT_BODY
    p5.font.size = Pt(10.5)
    p5.font.color.rgb = TEXT_MUTED

    # Bottom Row: 3 Presenter Cards
    pres_cards = [
        ("V. Ram Charan", "Act I: Core Hardware & System Bus (Slides 1–5)", ACCENT_CYAN),
        ("Vedhanth", "Act II: Memory & Storage Hierarchy (Slides 6–10)", ACCENT_INDIGO),
        ("Lochan", "Act III: Multiprocessing & Modern Systems (Slides 11–15)", ACCENT_EMERALD)
    ]
    p_w = 3.7
    for i, (name, act_info, col) in enumerate(pres_cards):
        px = 0.8 + i * 4.0
        add_card(s15, px, 5.8, p_w, 1.2, name, border_color=CARD_BORDER, bg_color=PILL_BG, accent_bar=col)
        tb = s15.shapes.add_textbox(Inches(px + 0.2), Inches(6.25), Inches(p_w - 0.4), Inches(0.6))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = act_info
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = col

    set_speaker_notes(
        s15,
        "To summarize our complete presentation: Ram Charan walked us through the core hardware and 3-part system bus; Vedhanth explained the storage hierarchy, locality of reference, and OS caching optimizations; and Lochan covered multiprocessing architectures, SMP vs AMP, multicore, clustering, and OS coordination challenges. On behalf of Ram Charan, Vedhanth, and Lochan, thank you for your attention. We are now open for questions.",
        "Comprehensive architectural synthesis covering core hardware, system bus, storage hierarchy, locality, multiprocessing, and OS coordination.",
        "End of presentation — open floor for Q&A.",
        "All three presenters step forward together, smile, and invite questions from the professor and audience."
    )

    out_pptx = os.path.join(ROOT_DIR, "computer_system_architecture.pptx")
    prs.save(out_pptx)
    out_pptx_clean = os.path.join(ROOT_DIR, "Computer_System_Architecture_OS.pptx")
    prs.save(out_pptx_clean)
    out_pptx_cinematic = os.path.join(ROOT_DIR, "Computer_System_Architecture_Cinematic.pptx")
    prs.save(out_pptx_cinematic)
    print(f"Successfully generated 15-slide Luxury presentation updated to exact user speech!")
    return out_pptx

if __name__ == "__main__":
    build_presentation()
